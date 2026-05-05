"""
월별 스터디 PPT 생성 스크립트

사용법:
  python generate_ppt.py --year-month 2604

동작:
  1. Drive YYMM/음성전사/ 에서 전사 텍스트 읽기
  2. Confluence 회의록 (월별 하위 페이지) 읽기
  3. Slack 스터디 인증 사진 다운로드
  4. python-pptx 로 PPT 생성
  5. Drive YYMM/ 에 PPT 업로드
"""

import argparse
import os
import re
import tempfile
from io import BytesIO
from pathlib import Path

import requests
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

from config import (
    DRIVE_ROOT_FOLDER_ID,
    CONFLUENCE_BASE_URL,
    CONFLUENCE_EMAIL,
    CONFLUENCE_API_TOKEN,
    CONFLUENCE_ROOT_PAGE_ID,
    SLACK_BOT_TOKEN,
    SLACK_CHANNEL_ID,
    TRANSCRIPT_SUBFOLDER,
)
from drive_utils import (
    get_drive_service,
    find_or_create_folder,
    list_files_in_folder,
    download_file,
    upload_file,
)
from slack_utils import fetch_study_photos

# ── 색상 팔레트 ──────────────────────────────────────────────
BRAND_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BRAND_BLUE = RGBColor(0x16, 0x21, 0x3E)
BRAND_ACCENT = RGBColor(0x0F, 0x3A, 0x71)
BRAND_LIGHT = RGBColor(0xE0, 0xE8, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Confluence 헬퍼 ──────────────────────────────────────────
def _cf_auth():
    from requests.auth import HTTPBasicAuth
    return HTTPBasicAuth(CONFLUENCE_EMAIL, CONFLUENCE_API_TOKEN)


def get_monthly_pages(year_month: str) -> list[dict]:
    """YYMM에 해당하는 월의 회의록 하위 페이지 목록 반환"""
    yy, mm = int(year_month[:2]) + 2000, int(year_month[2:])
    url = f"{CONFLUENCE_BASE_URL}/wiki/rest/api/content/{CONFLUENCE_ROOT_PAGE_ID}/child/page"
    resp = requests.get(url, auth=_cf_auth(), params={"limit": 100})
    resp.raise_for_status()
    pages = []
    for p in resp.json().get("results", []):
        title = p.get("title", "")
        # 제목이 'M/D' 또는 'M/DD' 형식인 경우만 수집 (해당 월)
        m = re.match(r"^(\d{1,2})/(\d{1,2})$", title)
        if m and int(m.group(1)) == mm:
            pages.append({"id": p["id"], "title": title})
    pages.sort(key=lambda x: int(x["title"].split("/")[1]))
    return pages


def get_page_content(page_id: str) -> dict:
    """회의록 페이지의 표와 insights 를 파싱해 반환"""
    url = f"{CONFLUENCE_BASE_URL}/wiki/rest/api/content/{page_id}"
    resp = requests.get(url, auth=_cf_auth(), params={"expand": "body.storage"})
    resp.raise_for_status()
    data = resp.json()
    html = data.get("body", {}).get("storage", {}).get("value", "")

    # 간단 파싱: <td> 추출 → 담당자/계획/한 것 구성
    rows = []
    for tr in re.findall(r"<tr>(.*?)</tr>", html, re.DOTALL):
        cells = re.findall(r"<td[^>]*>(.*?)</td>", tr, re.DOTALL)
        clean = [re.sub(r"<[^>]+>", "", c).strip() for c in cells]
        if len(clean) >= 3 and clean[0]:
            rows.append(clean[:3])

    # insights 섹션 추출
    insights_match = re.search(
        r"<h1[^>]*>insights</h1>(.*?)(?=<h[12]|$)", html, re.DOTALL | re.IGNORECASE
    )
    insights_html = insights_match.group(1) if insights_match else ""
    # <li> 항목 추출
    bullets = re.findall(r"<li[^>]*>(.*?)</li>", insights_html, re.DOTALL)
    insights = [re.sub(r"<[^>]+>", "", b).strip() for b in bullets if b.strip()]

    return {"rows": rows, "insights": insights}


# ── Drive 헬퍼 ───────────────────────────────────────────────
def _find_subfolder(service, parent_id, name):
    query = (
        f"name='{name}' and mimeType='application/vnd.google-apps.folder'"
        f" and '{parent_id}' in parents and trashed=false"
    )
    results = service.files().list(q=query, fields="files(id)").execute()
    files = results.get("files", [])
    if not files:
        return None
    return files[0]["id"]


def get_transcripts(service, year_month: str) -> dict[str, str]:
    """Drive YYMM/음성전사/ 에서 txt 파일을 읽어 {파일명: 내용} 반환"""
    month_id = _find_subfolder(service, DRIVE_ROOT_FOLDER_ID, year_month)
    if not month_id:
        return {}
    trans_id = _find_subfolder(service, month_id, TRANSCRIPT_SUBFOLDER)
    if not trans_id:
        return {}

    result = {}
    for f in list_files_in_folder(service, trans_id):
        if f["name"].endswith(".txt"):
            with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp:
                download_file(service, f["id"], tmp.name)
                with open(tmp.name, encoding="utf-8") as fp:
                    result[f["name"]] = fp.read()
                os.unlink(tmp.name)
    return result


# ── PPT 빌더 ────────────────────────────────────────────────
def _set_background(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def _add_rect(slide, left, top, width, height, fill_color: RGBColor, line=False):
    shape = slide.shapes.add_shape(
        1, left, top, width, height  # MSO_SHAPE_TYPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    return shape


def build_title_slide(prs: Presentation, year_month: str):
    yy, mm = int(year_month[:2]) + 2000, int(year_month[2:])
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_background(slide, BRAND_DARK)

    # 상단 액센트 바
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), YELLOW)

    # 타이틀
    _add_textbox(
        slide, "BCS 스터디",
        Inches(1.5), Inches(1.8), Inches(10), Inches(1),
        font_size=44, bold=True, color=YELLOW, align=PP_ALIGN.CENTER
    )
    _add_textbox(
        slide, f"월간 스터디 요약 | {yy}년 {mm}월",
        Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
        font_size=28, bold=False, color=WHITE, align=PP_ALIGN.CENTER
    )
    _add_textbox(
        slide, "AI Agent 활용 능력 향상 스터디",
        Inches(1.5), Inches(3.9), Inches(10), Inches(0.6),
        font_size=18, color=BRAND_LIGHT, align=PP_ALIGN.CENTER
    )

    # 하단 바
    _add_rect(slide, 0, Inches(7.0), SLIDE_W, Inches(0.5), BRAND_ACCENT)
    _add_textbox(
        slide, "매주 화요일 점심 · 2026.02.24 ~ 2026.06.30",
        Inches(0.5), Inches(7.05), Inches(12), Inches(0.4),
        font_size=12, color=BRAND_LIGHT, align=PP_ALIGN.CENTER
    )


def build_overview_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, BRAND_BLUE)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), YELLOW)
    _add_textbox(slide, "스터디 프로젝트 현황", Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                 font_size=28, bold=True, color=YELLOW)

    members = [
        ("이수영", "웹 게임 만들기 / 주식 추천 앱"),
        ("유민우", "냉장고 관리 앱 / 오이 포커 (Godot 멀티플레이)"),
        ("이준영", "Chiptune MIDI 장치 / 3D 프린팅 커뮤니티"),
        ("권종범", "금칙어 관리 서비스 / Klaf 에이전트 빌더"),
        ("이선용", "[웹] 프리셀 / 포모도로"),
    ]

    row_h = Inches(0.9)
    for i, (name, project) in enumerate(members):
        y = Inches(1.2) + i * row_h
        bg_color = BRAND_ACCENT if i % 2 == 0 else BRAND_BLUE
        _add_rect(slide, Inches(0.5), y, Inches(12.3), row_h - Inches(0.05), bg_color)
        _add_textbox(slide, name, Inches(0.7), y + Inches(0.1), Inches(2), Inches(0.7),
                     font_size=16, bold=True, color=YELLOW)
        _add_textbox(slide, project, Inches(2.9), y + Inches(0.1), Inches(9.5), Inches(0.7),
                     font_size=15, color=WHITE)


def build_meeting_slide(prs: Presentation, title: str, content: dict, transcript: str = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, BRAND_BLUE)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), YELLOW)

    # 날짜 헤더
    _add_textbox(slide, f"📅 {title} 회의록", Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                 font_size=26, bold=True, color=YELLOW)

    # 진행 내용 표 (왼쪽)
    _add_textbox(slide, "▪ 진행 내용", Inches(0.5), Inches(1.0), Inches(7.5), Inches(0.4),
                 font_size=14, bold=True, color=BRAND_LIGHT)

    row_h = Inches(0.55)
    visible_rows = content.get("rows", [])[:5]
    for i, row in enumerate(visible_rows):
        y = Inches(1.45) + i * row_h
        if i % 2 == 0:
            _add_rect(slide, Inches(0.5), y, Inches(7.5), row_h - Inches(0.03), BRAND_ACCENT)
        name = row[0] if row else ""
        done = row[2] if len(row) > 2 else ""
        # 멘션 태그 제거
        name = re.sub(r"@\S+", "", name).strip()
        done_short = done[:60] + "…" if len(done) > 60 else done
        _add_textbox(slide, name, Inches(0.6), y + Inches(0.05), Inches(1.8), Inches(0.5),
                     font_size=12, bold=True, color=YELLOW)
        _add_textbox(slide, done_short, Inches(2.5), y + Inches(0.05), Inches(5.3), Inches(0.5),
                     font_size=11, color=WHITE)

    # Insights (오른쪽)
    insights = content.get("insights", [])
    if insights:
        _add_textbox(slide, "💡 Insights", Inches(8.3), Inches(1.0), Inches(4.7), Inches(0.4),
                     font_size=14, bold=True, color=BRAND_LIGHT)
        ins_y = Inches(1.45)
        for ins in insights[:6]:
            short = ins[:70] + "…" if len(ins) > 70 else ins
            _add_textbox(slide, f"• {short}", Inches(8.3), ins_y, Inches(4.7), Inches(0.7),
                         font_size=11, color=WHITE, wrap=True)
            ins_y += Inches(0.75)

    # 음성전사 미리보기 (있을 때)
    if transcript:
        preview = transcript[:200].replace("\n", " ") + "…"
        _add_rect(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.85), BRAND_DARK)
        _add_textbox(slide, f"🎙 음성전사: {preview}", Inches(0.6), Inches(6.55),
                     Inches(12.1), Inches(0.75), font_size=10, color=BRAND_LIGHT)


def build_transcript_summary_slide(prs: Presentation, transcripts: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, BRAND_DARK)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), YELLOW)
    _add_textbox(slide, "🎙 음성 전사 요약", Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                 font_size=26, bold=True, color=YELLOW)

    if not transcripts:
        _add_textbox(
            slide,
            "음성 전사 파일이 아직 업로드되지 않았습니다.\n"
            "transcribe.py 스크립트를 실행한 후 PPT를 재생성하세요.",
            Inches(1), Inches(2.5), Inches(11), Inches(2),
            font_size=18, color=BRAND_LIGHT, align=PP_ALIGN.CENTER
        )
        return

    y = Inches(1.1)
    for fname, text in list(transcripts.items())[:3]:
        session = Path(fname).stem
        _add_textbox(slide, f"📁 {session}", Inches(0.5), y, Inches(12), Inches(0.4),
                     font_size=14, bold=True, color=YELLOW)
        y += Inches(0.4)
        preview = text[:300].replace("\n", " ") + ("…" if len(text) > 300 else "")
        _add_textbox(slide, preview, Inches(0.7), y, Inches(12), Inches(1.0),
                     font_size=11, color=WHITE, wrap=True)
        y += Inches(1.1)


def build_photo_slide(prs: Presentation, photo_path: str, caption: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, BRAND_DARK)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), YELLOW)
    _add_textbox(slide, f"📸 스터디 인증 사진 {caption}", Inches(0.5), Inches(0.1),
                 Inches(12), Inches(0.6), font_size=22, bold=True, color=YELLOW)

    try:
        img = Image.open(photo_path)
        iw, ih = img.size
        # 슬라이드에 맞게 비율 유지
        max_w = Inches(12)
        max_h = Inches(6.5)
        ratio = min(max_w / iw, max_h / ih)
        nw, nh = int(iw * ratio), int(ih * ratio)
        left = (SLIDE_W - nw) // 2
        top = Inches(0.8) + (max_h - nh) // 2
        slide.shapes.add_picture(photo_path, left, top, nw, nh)
    except Exception as e:
        _add_textbox(slide, f"[이미지 로드 실패: {e}]", Inches(1), Inches(3),
                     Inches(11), Inches(1), font_size=14, color=GRAY)


# ── 메인 ────────────────────────────────────────────────────
def run(year_month: str):
    service = get_drive_service()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("1/5 타이틀 슬라이드 생성...")
    build_title_slide(prs, year_month)

    print("2/5 프로젝트 개요 슬라이드 생성...")
    build_overview_slide(prs)

    print("3/5 Confluence 회의록 슬라이드 생성...")
    pages = get_monthly_pages(year_month)
    transcripts = get_transcripts(service, year_month)

    for page in pages:
        print(f"   {page['title']} ...")
        content = get_page_content(page["id"])
        # 날짜 매칭 전사 파일 찾기
        day = page["title"].split("/")[1].zfill(2)
        mm = page["title"].split("/")[0].zfill(2)
        yy = year_month[:2]
        trans_key_hint = f"{yy[0:2]}-{mm}-{day}"
        matching_trans = next(
            (v for k, v in transcripts.items() if mm in k and day in k), None
        )
        build_meeting_slide(prs, page["title"], content, matching_trans)

    print("4/5 음성전사 요약 슬라이드 생성...")
    build_transcript_summary_slide(prs, transcripts)

    print("5/5 스터디 인증 사진 슬라이드 생성...")
    with tempfile.TemporaryDirectory() as tmpdir:
        photo_paths = fetch_study_photos(year_month, tmpdir)
        for i, photo_path in enumerate(photo_paths):
            build_photo_slide(prs, photo_path, f"({i + 1}/{len(photo_paths)})")

        # PPT 저장 및 업로드
        ppt_name = f"바코스_스터디_{year_month}.pptx"
        local_ppt = os.path.join(tmpdir, ppt_name)
        prs.save(local_ppt)
        print(f"PPT 저장: {local_ppt}")

        month_folder_id = _find_subfolder(service, DRIVE_ROOT_FOLDER_ID, year_month)
        if not month_folder_id:
            from drive_utils import find_or_create_folder
            month_folder_id = find_or_create_folder(service, year_month, DRIVE_ROOT_FOLDER_ID)

        file_id = upload_file(
            service, local_ppt, ppt_name, month_folder_id,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        print(f"✅ 업로드 완료! Drive 파일 ID: {file_id}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="월별 스터디 PPT 생성")
    parser.add_argument("--year-month", required=True, help="YYMM 형식 (예: 2604)")
    args = parser.parse_args()
    run(args.year_month)
