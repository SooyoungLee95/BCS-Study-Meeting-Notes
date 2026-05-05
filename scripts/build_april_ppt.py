"""
2604 (2026년 4월) 스터디 PPT 로컬 생성 스크립트
Confluence 회의록 데이터를 직접 사용 (API 인증 불필요)
"""

import os, sys, re, tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BRAND_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
BRAND_BLUE   = RGBColor(0x16, 0x21, 0x3E)
BRAND_ACCENT = RGBColor(0x0F, 0x3A, 0x71)
BRAND_LIGHT  = RGBColor(0xE0, 0xE8, 0xF5)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY         = RGBColor(0x88, 0x88, 0x88)
YELLOW       = RGBColor(0xFF, 0xD7, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── 수집한 Confluence 회의록 데이터 ──────────────────────────
MEETINGS = [
    {
        "date": "4/7",
        "rows": [
            ("이수영", "anthropic courses 수강 (완강) / MCP server client 구현"),
            ("유민우", "오이 포커 1대1 완성"),
            ("이준영", "Claude Code 원격 서버 구축 완료 / 3D 모델 에디터 crash reporter / 모서리 처리"),
            ("권종범", "Klaf 에이전트 빌더 플러그인 구현 / Marklas transform 기능 추가"),
            ("이선용", "포모도로 고도화 중 토큰 소진"),
        ],
        "insights": [
            "프롬프트 사용량 줄이는 방법 시도 중",
            "클로드가 이미지 파일 디버깅은 잘 못함 → HTML 코드 문서로 변환하면 이해도 향상",
            "준영님 사업 모델 공유",
            "종범님: 에이전트별 스킬 분리 효과 (메인 컨텍스트 절약, 필요 스킬만 로드)",
            "Klaf 개발 이유: Claude Code가 에이전트를 스킬처럼 쓰는 경향 → 역할에 맞는 에이전트 개발 유도",
        ],
    },
    {
        "date": "4/14",
        "rows": [
            ("이수영", "없음"),
            ("유민우", "오이 포커 멀티플레이 구현 / Godot 엔진 공부"),
            ("이준영", "Claude Code 원격 서버 구축 / crash reporter / 모서리 처리 (휴가 예정)"),
            ("권종범", "craken 개인용 Claude Code 마켓플레이스 추가 / klaf 고도화"),
            ("이선용", "프리셀 완료 / schedule 이용 자동 기능 개발"),
        ],
        "insights": [
            "Claude schedule로 GitHub 이슈를 긁어와서 자동 개발",
            "토큰 초기화 주기(1시, 6시)마다 자동 실행",
            "1번 배치에 토큰 소진 → 다음 배치에 처음부터 시작 이슈 발생",
            "작업 현황 기록 방법 탐색 중",
        ],
    },
    {
        "date": "4/21",
        "rows": [
            ("이수영", "회의록 작성 워크플로우 구성"),
            ("유민우", "오이 포커 멀티플레이 구현 / Claude Design으로 화면 구성"),
            ("이준영", "3D 프린팅 모델 에디터 crash reporter / 커뮤니티 백엔드 계정 체계 / Outbox 패턴 구현"),
            ("권종범", "Klaf 고도화"),
            ("이선용", "하네스 공부"),
        ],
        "insights": [
            "클로드 디자인 공유",
            "Opus 4.7 동작 방식 변경 공유",
            "Klaf: .claude 루트에 스킬 수백 개 시 메인 에이전트가 스킬 선택을 잘 못함",
            "에이전트별 스킬 계층화 방안 고민 (python agent, java agent별 전용 스킬)",
            "책임의 상속 구조 가능 여부 검토 (python 개발자 → django 개발자)",
            "에이전트 하네스 공부",
        ],
    },
    {
        "date": "4/29",
        "rows": [
            ("이수영", "회의록 작성 워크플로우 구성 / claude routine으로 시도"),
            ("유민우", "오이 포커 멀티플레이 재구현 / Godot MCP 외 다른 Skills 레포 활용 개발"),
            ("이준영", "3D 프린팅 커뮤니티 백엔드 메인 로직 구현 / MCP 서버 구현 중"),
            ("권종범", "codi 구조 확인 / Claude Agent SDK + Slack 병합 구조"),
            ("이선용", "하네스 공부"),
        ],
        "insights": [
            "오이포커: 일대일→멀티 전환 계속 실패 → 아예 새로 시작 / 게임엔진 MCP 설계문서 작성",
            "mumcp: 필요 기능만 묶은 커스텀 MCP / 버스 배차+날씨로 퇴근 안내 루틴 목표",
            "claude simplify: 커밋 안된 변경사항, 현재 브랜치 조건 지정 가능",
            "codi repo 범용화 고민 (Slack 강결합 → 범용 구조로)",
            "하네스: 설정에 따라 성능 차이 큼 / 코드 에이전트는 실수를 전제 / 결과를 지시, 방법 X",
        ],
    },
]

# ── 헬퍼 ─────────────────────────────────────────────────────
def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# ── 슬라이드 빌더 ────────────────────────────────────────────
def title_slide(prs):
    s = new_slide(prs)
    bg(s, BRAND_DARK)
    rect(s, 0, 0, SLIDE_W, Inches(0.07), YELLOW)
    # 배경 장식
    rect(s, Inches(10.5), Inches(1.0), Inches(3.5), Inches(6.0), RGBColor(0x0A, 0x0A, 0x20))
    txt(s, "BCS 스터디", Inches(1.2), Inches(1.6), Inches(9.5), Inches(1.2),
        size=52, bold=True, color=YELLOW, align=PP_ALIGN.LEFT)
    txt(s, "월간 스터디 요약", Inches(1.2), Inches(2.9), Inches(9.5), Inches(0.9),
        size=32, color=WHITE, align=PP_ALIGN.LEFT)
    txt(s, "2026년 4월 | 2604", Inches(1.2), Inches(3.8), Inches(9.5), Inches(0.7),
        size=22, color=BRAND_LIGHT, align=PP_ALIGN.LEFT)
    rect(s, Inches(1.2), Inches(4.6), Inches(4.0), Inches(0.06), YELLOW)
    txt(s, "AI Agent 활용 능력 향상 스터디",
        Inches(1.2), Inches(4.75), Inches(9.5), Inches(0.5),
        size=15, color=BRAND_LIGHT)
    txt(s, "매주 화요일 점심 · 2026.02.24 ~ 2026.06.30",
        Inches(1.2), Inches(5.35), Inches(9.5), Inches(0.4),
        size=13, color=GRAY)
    # 우측 장식 숫자
    txt(s, "4", Inches(10.8), Inches(1.5), Inches(2.5), Inches(4.0),
        size=160, bold=True, color=RGBColor(0x22, 0x22, 0x44), align=PP_ALIGN.CENTER)
    txt(s, "월", Inches(10.8), Inches(5.2), Inches(2.5), Inches(0.8),
        size=26, color=RGBColor(0x44, 0x44, 0x66), align=PP_ALIGN.CENTER)
    rect(s, 0, Inches(7.0), SLIDE_W, Inches(0.5), BRAND_ACCENT)
    txt(s, "스터디원: 이수영 · 유민우 · 이준영 · 권종범 · 이선용",
        Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
        size=13, color=BRAND_LIGHT, align=PP_ALIGN.CENTER)


def overview_slide(prs):
    s = new_slide(prs)
    bg(s, BRAND_BLUE)
    rect(s, 0, 0, SLIDE_W, Inches(0.07), YELLOW)
    txt(s, "스터디 프로젝트 현황", Inches(0.6), Inches(0.2), Inches(11), Inches(0.7),
        size=30, bold=True, color=YELLOW)
    txt(s, "스터디 목표: 개인별 프로젝트 최소 1개 완성 | AI Agent 활용 능력 향상",
        Inches(0.6), Inches(0.9), Inches(12), Inches(0.4), size=13, color=BRAND_LIGHT)
    rect(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.03), YELLOW)

    members = [
        ("이수영", "웹 게임 만들기 / 주식 추천 앱 / 회의록 자동화 워크플로우"),
        ("유민우", "냉장고 관리 앱 / 오이 포커 (Godot 멀티플레이어)"),
        ("이준영", "Chiptune MIDI 장치 / 3D 프린팅 모델 공유 커뮤니티"),
        ("권종범", "금칙어 관리 서비스 / Klaf 에이전트 빌더 / Atlassian MCP"),
        ("이선용", "[웹] 프리셀 / 포모도로 / 하네스 공부"),
    ]
    row_h = Inches(1.0)
    icons = ["💻", "🎮", "🖨", "🤖", "🃏"]
    for i, (name, proj) in enumerate(members):
        y = Inches(1.45) + i * row_h
        c = BRAND_ACCENT if i % 2 == 0 else RGBColor(0x12, 0x1A, 0x35)
        rect(s, Inches(0.5), y, Inches(12.3), row_h - Inches(0.06), c)
        txt(s, icons[i], Inches(0.65), y + Inches(0.2), Inches(0.5), Inches(0.6),
            size=20, color=WHITE)
        txt(s, name, Inches(1.25), y + Inches(0.15), Inches(2.2), Inches(0.7),
            size=18, bold=True, color=YELLOW)
        txt(s, proj, Inches(3.6), y + Inches(0.15), Inches(8.9), Inches(0.7),
            size=15, color=WHITE)


def meeting_slide(prs, meeting: dict, idx: int, total: int):
    s = new_slide(prs)
    bg(s, BRAND_BLUE)
    rect(s, 0, 0, SLIDE_W, Inches(0.07), YELLOW)

    date_txt = meeting["date"]
    rect(s, Inches(0.5), Inches(0.12), Inches(1.8), Inches(0.72), BRAND_ACCENT)
    txt(s, date_txt, Inches(0.55), Inches(0.15), Inches(1.7), Inches(0.66),
        size=26, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    txt(s, "2026년 4월 회의록", Inches(2.4), Inches(0.18), Inches(8), Inches(0.55),
        size=24, bold=True, color=WHITE)
    txt(s, f"({idx}/{total})", Inches(11.5), Inches(0.18), Inches(1.3), Inches(0.55),
        size=14, color=GRAY, align=PP_ALIGN.RIGHT)
    rect(s, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.03), RGBColor(0x2A, 0x3A, 0x5A))

    # 왼쪽: 진행 내용 표
    rect(s, Inches(0.5), Inches(1.0), Inches(7.7), Inches(0.36), BRAND_ACCENT)
    txt(s, "▪ 이번주 진행 내용", Inches(0.6), Inches(1.03), Inches(7.5), Inches(0.3),
        size=13, bold=True, color=WHITE)

    row_h = Inches(0.54)
    for i, (name, done) in enumerate(meeting["rows"]):
        y = Inches(1.38) + i * row_h
        c = RGBColor(0x10, 0x2A, 0x55) if i % 2 == 0 else RGBColor(0x12, 0x1E, 0x40)
        rect(s, Inches(0.5), y, Inches(7.7), row_h - Inches(0.04), c)
        txt(s, name, Inches(0.62), y + Inches(0.07), Inches(1.7), Inches(0.44),
            size=12, bold=True, color=YELLOW)
        short = done[:58] + "…" if len(done) > 58 else done
        txt(s, short, Inches(2.45), y + Inches(0.07), Inches(5.6), Inches(0.44),
            size=11, color=WHITE)

    # 오른쪽: Insights
    rect(s, Inches(8.4), Inches(1.0), Inches(4.4), Inches(0.36), RGBColor(0x1A, 0x3A, 0x60))
    txt(s, "💡 Insights", Inches(8.5), Inches(1.03), Inches(4.2), Inches(0.3),
        size=13, bold=True, color=YELLOW)

    ins_y = Inches(1.43)
    for ins in meeting["insights"][:6]:
        short = ins[:64] + "…" if len(ins) > 64 else ins
        rect(s, Inches(8.4), ins_y, Inches(4.4), Inches(0.7), RGBColor(0x0D, 0x20, 0x3D))
        txt(s, f"• {short}", Inches(8.5), ins_y + Inches(0.06), Inches(4.2), Inches(0.62),
            size=10, color=BRAND_LIGHT, wrap=True)
        ins_y += Inches(0.74)

    # 하단 구분선
    rect(s, 0, Inches(7.05), SLIDE_W, Inches(0.45), BRAND_DARK)
    txt(s, f"BCS 스터디 · 2026년 4월 · {date_txt}",
        Inches(0.5), Inches(7.08), Inches(12.3), Inches(0.35),
        size=11, color=GRAY, align=PP_ALIGN.CENTER)


def transcript_slide(prs):
    s = new_slide(prs)
    bg(s, BRAND_DARK)
    rect(s, 0, 0, SLIDE_W, Inches(0.07), YELLOW)
    txt(s, "🎙 음성 전사 요약", Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
        size=28, bold=True, color=YELLOW)

    sessions = [
        ("4/7 회의 (바코스 4-7.m4a)", "4월 7일 회의 녹취록"),
        ("4/21 회의 (바코스 4-21.m4a)", "4월 21일 회의 녹취록"),
        ("4/29 회의 (바코스 4-29.m4a)", "4월 29일 회의 녹취록"),
    ]
    note = (
        "※ 음성 전사는 Google Drive 바코스/models/ggml-large-v3.bin (Whisper large-v3)\n"
        "   모델을 사용하여 transcribe.py 스크립트로 생성됩니다.\n"
        "   전사 파일은 Drive 바코스/2604/음성전사/ 폴더에 업로드됩니다."
    )

    rect(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.0), BRAND_ACCENT)
    txt(s, note, Inches(0.7), Inches(1.05), Inches(11.9), Inches(0.9),
        size=12, color=BRAND_LIGHT, wrap=True)

    for i, (session, desc) in enumerate(sessions):
        y = Inches(2.2) + i * Inches(1.5)
        rect(s, Inches(0.5), y, Inches(12.3), Inches(1.4), RGBColor(0x0F, 0x1A, 0x30))
        rect(s, Inches(0.5), y, Inches(0.08), Inches(1.4), YELLOW)
        txt(s, f"📁 {session}", Inches(0.72), y + Inches(0.1), Inches(11.8), Inches(0.45),
            size=14, bold=True, color=YELLOW)
        txt(s, f"   ↳ pipeline.py --year-month 2604 실행 후 음성전사 내용이 여기에 표시됩니다",
            Inches(0.72), y + Inches(0.6), Inches(11.8), Inches(0.65),
            size=11, color=BRAND_LIGHT)

    rect(s, 0, Inches(7.05), SLIDE_W, Inches(0.45), BRAND_ACCENT)
    txt(s, "실행: cd scripts && python pipeline.py --year-month 2604",
        Inches(0.5), Inches(7.08), Inches(12.3), Inches(0.35),
        size=12, color=WHITE, align=PP_ALIGN.CENTER)


def photo_slide(prs, date_str: str, image_path: str, idx: int, total: int):
    """Drive 인증사진을 슬라이드에 임베드"""
    from PIL import Image as PILImage
    s = new_slide(prs)
    bg(s, BRAND_DARK)
    rect(s, 0, 0, SLIDE_W, Inches(0.07), YELLOW)
    txt(s, f"📸 스터디 인증 사진 ({idx}/{total}) · {date_str}",
        Inches(0.5), Inches(0.12), Inches(12), Inches(0.65),
        size=24, bold=True, color=YELLOW)

    img = PILImage.open(image_path)
    iw, ih = img.size
    avail_w = Inches(11.5)
    avail_h = Inches(5.8)
    aw_emu = int(avail_w)
    ah_emu = int(avail_h)
    img_ratio = iw / ih
    box_ratio = aw_emu / ah_emu
    if img_ratio > box_ratio:
        final_w = aw_emu
        final_h = int(aw_emu / img_ratio)
    else:
        final_h = ah_emu
        final_w = int(ah_emu * img_ratio)

    left = (int(SLIDE_W) - final_w) // 2
    top = Inches(1.0)
    s.shapes.add_picture(image_path, left, top, width=final_w, height=final_h)

    rect(s, 0, Inches(7.05), SLIDE_W, Inches(0.45), BRAND_ACCENT)
    txt(s, f"BCS 스터디 인증 · 2026년 4월 · {date_str}",
        Inches(0.5), Inches(7.08), Inches(12.3), Inches(0.35),
        size=11, color=BRAND_LIGHT, align=PP_ALIGN.CENTER)


def summary_slide(prs):
    s = new_slide(prs)
    bg(s, BRAND_DARK)
    rect(s, 0, 0, SLIDE_W, Inches(0.07), YELLOW)
    txt(s, "2026년 4월 스터디 요약", Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
        size=30, bold=True, color=YELLOW)

    highlights = [
        ("🎮 유민우", "오이 포커 Godot 멀티플레이어 전면 재구현 / Claude Design 화면 구성"),
        ("🖨 이준영", "3D 프린팅 커뮤니티 백엔드 Outbox 패턴 + 보안 감사 / MCP 서버 개발 착수"),
        ("🤖 권종범", "Klaf 에이전트 계층화 설계 / codi Slack 범용화 / Claude Agent SDK+Slack"),
        ("🃏 이선용", "프리셀 완성 → UI 개선 / 하네스(Harness) 심층 공부"),
        ("💻 이수영", "Anthropic API 과정 완강 / MCP client-server 구현 / 회의록 자동화 시작"),
    ]
    for i, (name, desc) in enumerate(highlights):
        y = Inches(1.05) + i * Inches(1.1)
        c = BRAND_ACCENT if i % 2 == 0 else RGBColor(0x10, 0x1E, 0x3A)
        rect(s, Inches(0.5), y, Inches(12.3), Inches(1.0), c)
        rect(s, Inches(0.5), y, Inches(0.1), Inches(1.0), YELLOW)
        txt(s, name, Inches(0.75), y + Inches(0.1), Inches(2.5), Inches(0.8),
            size=15, bold=True, color=YELLOW)
        txt(s, desc, Inches(3.4), y + Inches(0.1), Inches(9.2), Inches(0.8),
            size=13, color=WHITE, wrap=True)

    rect(s, 0, Inches(7.05), SLIDE_W, Inches(0.45), BRAND_ACCENT)
    txt(s, "다음 스터디: 매주 화요일 점심 · AI Agent 활용 능력 향상",
        Inches(0.5), Inches(7.08), Inches(12.3), Inches(0.35),
        size=12, color=WHITE, align=PP_ALIGN.CENTER)


# ── 메인 ────────────────────────────────────────────────────
def build_ppt(output_path: str):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("1/7 타이틀 슬라이드...")
    title_slide(prs)

    print("2/7 프로젝트 개요 슬라이드...")
    overview_slide(prs)

    print("3/7 회의록 슬라이드 (4회) ...")
    for i, m in enumerate(MEETINGS):
        meeting_slide(prs, m, i + 1, len(MEETINGS))

    print("4/7 월간 요약 슬라이드...")
    summary_slide(prs)

    print("5/7 음성전사 슬라이드...")
    transcript_slide(prs)

    photo_dir = os.environ.get("BCS_PHOTO_DIR", "/tmp/bcs/photos")
    photo_files = []
    if os.path.isdir(photo_dir):
        import glob
        photo_files = sorted(
            glob.glob(os.path.join(photo_dir, "*.png")) +
            glob.glob(os.path.join(photo_dir, "*.jpg"))
        )

    if photo_files:
        for idx, path in enumerate(photo_files, 1):
            stem = Path(path).stem
            label = stem
            if stem.isdigit() and len(stem) == 4:
                label = f"4월 {int(stem[2:])}일"
            print(f"{5+idx}/{5+len(photo_files)} 스터디 인증 사진 슬라이드 ({label}) ...")
            photo_slide(prs, label, path, idx, len(photo_files))
    else:
        print("⚠ 인증사진 폴더 없음 (BCS_PHOTO_DIR), 사진 슬라이드 생략")

    prs.save(output_path)
    print(f"\n✅ PPT 저장 완료: {output_path}")
    return output_path


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "/tmp/바코스_스터디_2604.pptx"
    build_ppt(out)
